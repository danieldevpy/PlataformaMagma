# -*- coding: utf-8 -*-
"""
Renderizador de decks "Magma Cursos" — python-pptx + Pillow.

Data-driven: os slides são descritos como dados (o MESMO schema que o Studio
HTML exporta em deck.json). O mesmo arquivo serve para:

    python montar_powerpoint.py                  # renderiza o DECK embutido
    python montar_powerpoint.py deck.json        # renderiza um deck externo
    python montar_powerpoint.py deck.json out.pptx
    python montar_powerpoint.py --dump deck.json # grava o DECK embutido em JSON

Fluxo com o Studio:  Studio -> "Exportar deck.json" -> este script gera o .pptx.

Requisitos:  pip install python-pptx pillow

Schema de um slide (resumo):
    { "component": "content_card", "eyebrow": "...", "title": "...",
      "body": ["texto com **negrito**"], "footnote": "...", "stats": [...],
      "callout": {"kind":"alert","title":null,"desc":"..."} }

Componentes: cover, content_card, step_card, split, two_col, grid, grid3,
             do_dont, flow, hero.
"""
import os, re, sys, json, math

BASE = os.path.dirname(os.path.abspath(__file__))

def _find_img_dir():
    for cand in (os.path.join(BASE, "textos e imagens"),
                 os.path.join(BASE, "..", "textos e imagens")):
        if os.path.isdir(cand):
            return os.path.abspath(cand)
    return os.path.join(BASE, "textos e imagens")

IMG_DIR = _find_img_dir()
TMP_DIR = os.path.join(BASE, ".tmp_img")
DEFAULT_OUT = os.path.join(BASE, "lei_lucas_powerpoint_novo.pptx")
DATA = "17/07/2026"

# imports pesados só quando for realmente renderizar (permite --dump sem pptx)
def _import_pptx():
    global Presentation, Emu, Pt, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, Image
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from PIL import Image
    except ImportError as e:
        sys.exit("Faltam dependências. Rode:  pip install python-pptx pillow\n(%s)" % e)

# ---------------- tokens de cor (espelham app/css/tokens.css) -----------------
C = {
    "navy_deep": "101C38", "navy": "1B2A4D", "navy_soft": "24365E",
    "ink": "212A3D", "muted": "5B6476",
    "gold": "B8933F", "gold_light": "DCB96A", "gold_pale": "F0E3C4", "gold_deep": "8A6A1F",
    "red": "C8102E", "vida_blue": "1D4F91", "graphite": "232C3D",
    "paper": "FAF8F4", "sand": "ECEAE4", "white": "FFFFFF", "success": "2E8B57",
    "card_line": "E3DFD6",
    "alert_bg": "F8E6E9", "info_bg": "E9EFF9", "success_bg": "E5F1EA", "tip_bg": "F6EEDA",
}
def rgb(k):
    return RGBColor.from_string(C.get(k, k))

FONT_DISPLAY = "Arial"
FONT_BODY = "Arial"
DPI = 144.0
def PX(px):
    return Emu(int(px / DPI * 914400))
W, H = 1920, 1080

prs = None
BLANK = None
_imgc = [0]

def _new_presentation():
    global prs, BLANK
    prs = Presentation()
    prs.slide_width = PX(W); prs.slide_height = PX(H)
    BLANK = prs.slide_layouts[6]
    return prs

# ---------- limpeza de emoji (o design usa badges/ícones, não emoji) ----------
_EMOJI = re.compile(
    "[\U0001F000-\U0001FAFF\U00002600-\U000027BF\U0001F1E6-\U0001F1FF"
    "\U00002190-\U000021FF\U00002B00-\U00002BFF\U0000FE00-\U0000FE0F"
    "\U000023E9-\U000023FA\U00002B50\U0000203C\U00002049]", flags=re.UNICODE)
def clean(t):
    return re.sub(r"\s{2,}", " ", _EMOJI.sub("", t)).strip()

# ---------------- markdown inline compartilhado com o Studio ------------------
def parse_inline(text, base="ink", bold_color="navy"):
    """'A **Lei X** obriga' -> [('A ',False,base),('Lei X',True,bold_color),...]"""
    if not isinstance(text, str):
        return text  # já é uma lista de segmentos
    segs = []
    for i, part in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if part == "":
            continue
        bold = (i % 2 == 1)
        segs.append((part, bold, bold_color if bold else base))
    return segs or [(text, False, base)]

# ---------------- primitivas de desenho ---------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1, shape=None, round_val=None):
    if shape is None:
        shape = MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape, PX(x), PX(y), PX(w), PX(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if round_val is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = round_val
        except Exception: pass
    return sp

def set_spc(r, pts):
    r.font._rPr.set('spc', str(int(pts * 100)))

def run(p, text, size, color, bold=False, font=FONT_BODY, italic=False, spc=None):
    r = p.add_run(); r.text = clean(text)
    r.font.size = Pt(size); r.font.name = font; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = rgb(color)
    if spc is not None: set_spc(r, spc)
    return r

def _newframe(slide, x, y, w, h, anchor=None, wrap=True):
    if anchor is None:
        anchor = MSO_ANCHOR.TOP
    tb = slide.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tb, tf

def textbox(slide, x, y, w, h, text, size, color, bold=False, align=None,
            anchor=None, font=FONT_BODY, line_spacing=1.15, spc=None):
    if align is None: align = PP_ALIGN.LEFT
    if anchor is None: anchor = MSO_ANCHOR.TOP
    tb, tf = _newframe(slide, x, y, w, h, anchor)
    paras = text if isinstance(text, list) else [text]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        segs = para if isinstance(para, list) else [(para, bold, color)]
        for seg in segs:
            r = run(p, seg[0], size, seg[2] if len(seg) > 2 else color,
                    seg[1] if len(seg) > 1 else bold, font)
            if spc is not None: set_spc(r, spc)
    return tb

def est_lines(text, width_px, font_pt, cpl_factor=0.80):
    cpl = max(8, int(width_px / (font_pt * 1.333 * cpl_factor)))
    lines, cur = 1, 0
    for w in text.split():
        if cur + len(w) + 1 <= cpl: cur += len(w) + 1
        else: lines += 1; cur = len(w) + 1
    return lines

def pitch(size, ls):
    return size * 1.333 * ls * 1.55

def est_block_h(blocks, width, base=15):
    h = 0
    for b in blocks:
        k = b[0]
        if k == "intro":   h += est_lines(b[1], width, base) * pitch(base, 1.4) + 16
        elif k == "h":     h += pitch(19, 1.1) + 18
        elif k == "p":     h += est_lines("".join(s[0] for s in b[1]), width, base) * pitch(base, 1.45) + 12
        elif k == "bullet":h += est_lines("  " + b[1], width, base) * pitch(base, 1.3) + 8
        elif k == "step":  h += pitch(19, 1.12) + 14 + est_lines(b[3], width, base) * pitch(base, 1.38) + 8
        elif k in ("warn", "hint"): h += est_lines(b[1], width, base) * pitch(base + 1, 1.3) + 14
    return h

def callout_h(desc, width, title=None):
    return 34 + (32 if title else 0) + est_lines(desc, width - 110, 17) * 30 + 14

def gold_rule(slide, x, y, w=120, h=6):
    rect(slide, x, y, w, h, fill="gold_light")

def star_of_life(slide, x, y, size, color="white"):
    bl, bt = size * 0.60, size * 0.15
    cx, cy = x + size / 2, y + size / 2
    for ang in (0, 60, 120):
        bar = rect(slide, cx - bl / 2, cy - bt / 2, bl, bt, fill=color,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.5)
        bar.rotation = ang

def logo(slide, x, y, size=60, on_dark=True):
    hexa = rect(slide, x, y, size, size, fill="red", line="graphite", line_w=1.5,
                shape=MSO_SHAPE.HEXAGON); hexa.rotation = 90
    star_of_life(slide, x, y, size)
    wc = "white" if on_dark else "navy"
    textbox(slide, x + size + 16, y + size * 0.08, 380, size * 0.55, "MAGMA",
            int(size * 0.42), wc, True, font=FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    cur = textbox(slide, x + size + 17, y + size * 0.58, 380, size * 0.40, "CURSOS",
            int(size * 0.20), "gold", True, font=FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    for p in cur.text_frame.paragraphs:
        for r in p.runs: set_spc(r, 3.5)

def cover_crop(src, target_w, target_h):
    os.makedirs(TMP_DIR, exist_ok=True)
    _imgc[0] += 1
    out = os.path.join(TMP_DIR, f"c{_imgc[0]}.jpg")
    im = Image.open(src).convert("RGB"); iw, ih = im.size
    tr, ir = target_w / target_h, iw / ih
    if ir > tr:
        nw = int(ih * tr); left = (iw - nw) // 2; im = im.crop((left, 0, left + nw, ih))
    else:
        nh = int(iw / tr); top = (ih - nh) // 2; im = im.crop((0, top, iw, top + nh))
    im.save(out, quality=90); return out

def _is_photo(src):
    try:
        with Image.open(src) as im: return min(im.size) >= 400
    except Exception: return False

def place_image(slide, name, x, y, w, h):
    src = os.path.join(IMG_DIR, name) if name else None
    if src and os.path.exists(src) and _is_photo(src):
        slide.shapes.add_picture(cover_crop(src, w, h), PX(x), PX(y), PX(w), PX(h))
    elif name is None:
        rect(slide, x, y, w, h, fill="sand", line="card_line", line_w=1.2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.04)
        textbox(slide, x, y + h/2 - 20, w, 40, "Imagem", 16, "muted",
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def footer(slide, n, dark=False):
    col = "sand" if dark else "muted"
    y = 1028
    textbox(slide, 90, y, 500, 34, "MAGMA Cursos", 10.5, col, True, font=FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, W/2-250, y, 500, 34, DATA, 10.5, col, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, W-590, y, 500, 34, f"Slide {n:02d}", 10.5, col, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def header(slide, eyebrow, title, dark=False, x=90, title_size=38):
    gold_rule(slide, x, 150)
    if eyebrow:
        textbox(slide, x, 182, 1650, 34, eyebrow.upper(), 14.5,
                "gold_light" if dark else "gold", True, font=FONT_DISPLAY, spc=2.4)
    textbox(slide, x-2, 212, 1740, 132, title, title_size, "white" if dark else "navy",
            True, font=FONT_DISPLAY, line_spacing=1.05)
    return 360

def draw_callout(slide, x, y, w, kind, title, desc):
    m = {"alert": ("red", "alert_bg", "red"), "info": ("navy", "info_bg", "navy"),
         "success": ("success", "success_bg", "success"), "tip": ("gold", "tip_bg", "gold_deep")}[kind]
    ch = callout_h(desc, w, title)
    rect(slide, x, y, w, ch, fill=m[1], shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.05)
    rect(slide, x, y, 12, ch, fill=m[0])
    tb, tf = _newframe(slide, x + 40, y + 22, w - 70, ch - 30)
    if title:
        p = tf.paragraphs[0]; p.line_spacing = 1.1; p.space_after = Pt(5)
        run(p, title, 18, m[2], True, FONT_DISPLAY)
        p2 = tf.add_paragraph(); p2.line_spacing = 1.35
        run(p2, desc, 17, "ink")
    else:
        p = tf.paragraphs[0]; p.line_spacing = 1.35
        run(p, desc, 17, "ink")
    return ch

def build_flow(slide, x, y, w, blocks, base=15):
    tb, tf = _newframe(slide, x, y, w, H - y - 90)
    started = False
    for b in blocks:
        k = b[0]
        p = tf.paragraphs[0] if not started else tf.add_paragraph(); started = True
        if k == "warn":
            p.line_spacing = 1.3; p.space_before = Pt(12)
            run(p, "!  ", base + 2, "red", True, FONT_DISPLAY); run(p, b[1], base + 1, "red", True)
        elif k == "hint":
            p.line_spacing = 1.3; p.space_before = Pt(12)
            run(p, "»  ", base + 2, "gold", True, FONT_DISPLAY); run(p, b[1], base + 1, "gold_deep", True)
        elif k == "intro":
            p.line_spacing = 1.4; p.space_after = Pt(12); run(p, b[1], base, "muted")
        elif k == "h":
            p.line_spacing = 1.1; p.space_before = Pt(10); p.space_after = Pt(5)
            run(p, b[1], 19, "navy", True, FONT_DISPLAY)
        elif k == "p":
            p.line_spacing = 1.45; p.space_after = Pt(7)
            for seg in b[1]: run(p, seg[0], base, seg[2] if len(seg) > 2 else "ink", seg[1] if len(seg) > 1 else False)
        elif k == "bullet":
            p.line_spacing = 1.28; p.space_after = Pt(5)
            run(p, "•  ", base, "gold", True); run(p, b[1], base, "ink")
        elif k == "step":
            p.line_spacing = 1.08; p.space_before = Pt(7)
            run(p, b[1] + "   ", 19, "gold", True, FONT_DISPLAY)
            run(p, b[2], 19, "navy", True, FONT_DISPLAY)
            p2 = tf.add_paragraph(); p2.line_spacing = 1.3; p2.space_after = Pt(3)
            run(p2, b[3], base, "ink")
    return tb

# ---------------- COMPONENTES -------------------------------------------------
def comp_cover(eyebrow, title, subtitle, image=None):
    s = add_slide(); split = 1058
    rect(s, 0, 0, split, H, fill="navy_deep"); rect(s, split, 0, W - split, H, fill="paper")
    place_image(s, image, split, 0, W - split, H)
    gold_rule(s, 90, 205); logo(s, 90, 265, 60, True)
    textbox(s, 90, 512, 860, 36, eyebrow, 15, "gold_light", True, font=FONT_DISPLAY, spc=2.2)
    textbox(s, 88, 540, 900, 312, title, 44, "white", True, font=FONT_DISPLAY,
            line_spacing=1.05, anchor=MSO_ANCHOR.BOTTOM)
    textbox(s, 90, 880, 860, 150, subtitle, 18, "sand", line_spacing=1.35)
    return s

def comp_content_card(n, eyebrow, title, body, footnote=None, stats=None, callout=None):
    s = add_slide(); rect(s, 0, 0, W, H, fill="paper")
    gold_rule(s, 90, 150)
    textbox(s, 90, 182, 900, 36, eyebrow.upper(), 15, "gold", True, font=FONT_DISPLAY, spc=2.4)
    rect(s, 90, 250, 104, 104, fill="gold_pale", shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.22)
    hx = rect(s, 112, 272, 60, 60, fill="red", line="graphite", line_w=1.2, shape=MSO_SHAPE.HEXAGON); hx.rotation = 90
    star_of_life(s, 112, 272, 60)
    textbox(s, 224, 250, 1560, 104, title, 44, "navy", True, font=FONT_DISPLAY,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    tb, tf = _newframe(s, 224, 408, 1500, 560)
    started = [False]
    def par(segs, size, ls, sb=0, sa=0, font=FONT_BODY):
        p = tf.add_paragraph() if started[0] else tf.paragraphs[0]; started[0] = True
        p.line_spacing = ls
        if sb: p.space_before = Pt(sb)
        if sa: p.space_after = Pt(sa)
        for seg in segs: run(p, seg[0], size, seg[2] if len(seg) > 2 else "ink", seg[1] if len(seg) > 1 else False, font)
    for para in (body if isinstance(body, list) else [body]):
        par(para if isinstance(para, list) else [(para, False, "ink")], 22, 1.45, sa=6)
    if stats:
        for st in stats:
            par([("▎  ", True, "gold"), (st, True, "navy")], 24, 1.15, sb=14, font=FONT_DISPLAY)
    if footnote:
        par([(footnote, False, "muted")], 19, 1.4, sb=22)
    if callout:
        ch = callout_h(callout[2], 1700, callout[1])
        draw_callout(s, 90, 990 - ch, 1700, callout[0], callout[1], callout[2])
    footer(s, n)
    return s

def _place_callout(s, x, w, ct, content_h, kind, title, desc):
    ch = callout_h(desc, w, title)
    y = min(ct + content_h + 22, 990 - ch)
    draw_callout(s, x, y, w, kind, title, desc)

def comp_step_card(n, eyebrow, title, steps, intro=None, callout=None, image=None, note=None, hint=None):
    s = add_slide(); rect(s, 0, 0, W, H, fill="paper")
    ct = header(s, eyebrow, title)
    left_w = 1010 if image else 1680
    blocks = []
    if intro: blocks.append(("intro", intro))
    for st in steps: blocks.append(("step", st[0], st[1], st[2]))
    if note: blocks.append(("warn", note))
    if hint: blocks.append(("hint", hint))
    build_flow(s, 90, ct, left_w, blocks)
    if image:
        place_image(s, image, 1140, ct, 690, min(H - ct - 110, 560))
    if callout:
        _place_callout(s, 90, 1700, ct, est_block_h(blocks, left_w), *callout)
    footer(s, n)
    return s

def comp_split(n, eyebrow, title, blocks, image, callout=None):
    s = add_slide(); rect(s, 0, 0, W, H, fill="paper")
    ct = header(s, eyebrow, title)
    build_flow(s, 90, ct, 980, blocks)
    place_image(s, image, 1120, ct, 710, min(H - ct - 110, 590))
    if callout:
        _place_callout(s, 90, 990, ct, est_block_h(blocks, 980), *callout)
    footer(s, n)
    return s

def comp_two_col(n, eyebrow, title, left_blocks, right_blocks, callout=None):
    s = add_slide(); rect(s, 0, 0, W, H, fill="paper")
    ct = header(s, eyebrow, title)
    cw = 840; xr = 90 + cw + 60
    build_flow(s, 90, ct, cw, left_blocks)
    build_flow(s, xr, ct, cw, right_blocks)
    if callout:
        ch = callout_h(callout[2], 1700, callout[1])
        draw_callout(s, 90, 990 - ch, 1700, callout[0], callout[1], callout[2])
    footer(s, n)
    return s

def _card(slide, x, y, w, h, label, desc, badge=None):
    rect(slide, x, y, w, h, fill="white", line="card_line", line_w=1.2,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.05)
    pad = 26
    if badge:
        rect(slide, x + pad, y + pad, 46, 46, fill="gold", shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.24)
        textbox(slide, x + pad, y + pad, 46, 46, badge, 18, "white", True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_DISPLAY)
        textbox(slide, x + pad + 60, y + pad, w - 2*pad - 60, 48, label, 19, "navy", True,
                font=FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        ty = y + pad + 60
    else:
        textbox(slide, x + pad, y + pad, w - 2*pad, 40, label, 20, "navy", True, font=FONT_DISPLAY)
        ty = y + pad + 44
    textbox(slide, x + pad, ty, w - 2*pad, h - (ty - y) - pad, desc, 15.5, "ink", line_spacing=1.32)

def comp_grid(n, eyebrow, title, cards, cols=2, intro=None, badges=None, callout=None, quote=None):
    s = add_slide(); rect(s, 0, 0, W, H, fill="paper")
    ct = header(s, eyebrow, title)
    top = ct
    if quote:
        qh = est_lines(quote, 1600, 22) * 32 + 46
        rect(s, 90, top, 1740, qh, fill="tip_bg", shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.04)
        rect(s, 90, top, 12, qh, fill="gold")
        textbox(s, 130, top + 20, 1670, qh - 30, "“" + quote + "”", 22, "navy", True,
                font=FONT_DISPLAY, line_spacing=1.3)
        top += qh + 26
    if intro:
        il = est_lines(intro, 1740, 18)
        textbox(s, 90, top, 1740, il * 30 + 20, intro, 18, "muted", line_spacing=1.4)
        top += il * 30 + 28
    bottom = 1000
    ch_co = 0
    if callout:
        ch_co = est_lines(callout[2], 1600, 17) * 27 + 60
        bottom -= ch_co + 22
    gap = 26
    rows = math.ceil(len(cards) / cols)
    cw = (1740 - (cols - 1) * gap) / cols
    chh = (bottom - top - (rows - 1) * gap) / rows
    for i, cd in enumerate(cards):
        r_, c_ = divmod(i, cols)
        x = 90 + c_ * (cw + gap); y = top + r_ * (chh + gap)
        _card(s, x, y, cw, chh, cd[0], cd[1], badge=(badges[i] if badges else None))
    if callout:
        draw_callout(s, 90, 1000 - ch_co, 1700, callout[0], callout[1], callout[2])
    footer(s, n)
    return s

def comp_grid3(n, eyebrow, title, columns):
    s = add_slide(); rect(s, 0, 0, W, H, fill="paper")
    ct = header(s, eyebrow, title)
    gap = 28; cols = 3; cw = (1740 - 2 * gap) / cols
    y = ct; ch = 1000 - ct
    for i, (head, bullets, img) in enumerate(columns):
        x = 90 + i * (cw + gap)
        rect(s, x, y, cw, ch, fill="white", line="card_line", line_w=1.2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.04)
        place_image(s, img, x + 18, y + 18, cw - 36, 190)
        textbox(s, x + 24, y + 224, cw - 48, 46, head, 20, "navy", True, font=FONT_DISPLAY)
        tb, tf = _newframe(s, x + 24, y + 280, cw - 48, ch - 300)
        for j, bl in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = 1.3; p.space_after = Pt(7)
            run(p, "•  ", 15.5, "gold", True); run(p, bl, 15.5, "ink")
    footer(s, n)
    return s

def comp_do_dont(n, eyebrow, title, do_items, dont_items, callout=None):
    s = add_slide(); rect(s, 0, 0, W, H, fill="paper")
    ct = header(s, eyebrow, title)
    gap = 30; cw = (1740 - gap) / 2
    bottom = 1000
    ch_co = 0
    if callout:
        ch_co = est_lines(callout[2], 1600, 17) * 27 + 60
        bottom -= ch_co + 20
    colh = bottom - ct
    for idx, (items, lab, col, bg, mark) in enumerate([
            (do_items, "FAÇA", "success", "success_bg", "✔"),
            (dont_items, "NÃO FAÇA", "red", "alert_bg", "✕")]):
        x = 90 + idx * (cw + gap)
        rect(s, x, ct, cw, colh, fill=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.03)
        rect(s, x, ct, cw, 60, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_val=0.03)
        textbox(s, x + 28, ct, cw - 40, 60, lab, 18, "white", True, font=FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE, spc=1.5)
        tb, tf = _newframe(s, x + 30, ct + 86, cw - 60, colh - 110)
        for j, it in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = 1.28; p.space_after = Pt(10)
            run(p, mark + "  ", 16, col, True); run(p, it, 16.5, "ink")
    if callout:
        draw_callout(s, 90, 1000 - ch_co, 1700, callout[0], callout[1], callout[2])
    footer(s, n)
    return s

def comp_flow(n, eyebrow, title, stages, caption):
    s = add_slide(); rect(s, 0, 0, W, H, fill="navy_deep")
    header(s, eyebrow, title, dark=True)
    ns = len(stages); cw = 500; y = 540; hh = 165
    span = 1740; step = (span - cw) / max(1, (ns - 1))
    for i, st in enumerate(stages):
        x = 70 + i * step
        last = (i == ns - 1)
        fill = "gold" if last else ("navy_soft" if i % 2 == 0 else "navy")
        ch = rect(s, x, y, cw, hh, fill=fill, shape=MSO_SHAPE.CHEVRON)
        ch.line.color.rgb = rgb("gold_light" if last else "navy_soft"); ch.line.width = Pt(1)
        tb, tf = _newframe(s, x + 60, y, cw - 90, hh, anchor=MSO_ANCHOR.MIDDLE)
        for j, ln in enumerate(st.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.05
            run(p, ln, 21, "navy" if last else "white", True, FONT_DISPLAY)
    textbox(s, 160, 770, 1600, 200, caption, 18, "sand", align=PP_ALIGN.CENTER, line_spacing=1.45)
    footer(s, n, dark=True)
    return s

def comp_hero(n, title, subtitle=None, eyebrow=None, contact=None, tagline=None, big=60, dark=True):
    s = add_slide(); rect(s, 0, 0, W, H, fill="navy_deep")
    rect(s, 0, 0, W, 360, fill="navy_soft")
    rect(s, 0, 300, W, H - 300, fill="navy_deep")
    gold_rule(s, 90, 300)
    if eyebrow:
        textbox(s, 90, 336, 1000, 34, eyebrow.upper(), 15, "gold_light", True, font=FONT_DISPLAY, spc=2.4)
    textbox(s, 88, 360, 1500, 380, title, big, "white", True, font=FONT_DISPLAY,
            line_spacing=1.03, anchor=MSO_ANCHOR.BOTTOM)
    yy = 770
    if subtitle:
        textbox(s, 90, yy, 1300, 130, subtitle, 20, "sand", line_spacing=1.4); yy += 130
    if contact:
        for lead, desc in contact:
            textbox(s, 90, yy, 1300, 40, [[(lead + "  ", True, "gold_light"), (desc, False, "sand")]],
                    18, "sand", font=FONT_DISPLAY); yy += 46
        yy += 10
    if tagline:
        textbox(s, 90, yy, 1300, 60, tagline, 22, "gold_light", True, font=FONT_DISPLAY)
    logo(s, 90, 940, 54, True)
    textbox(s, W-590, 1028, 500, 34, f"Slide {n:02d}", 10.5, "sand", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return s

# ==================== ponte JSON -> componentes ==============================
def _callout(c):
    if not c: return None
    if isinstance(c, dict):
        return (c["kind"], c.get("title"), c["desc"])
    return tuple(c)

def _body(body):
    out = []
    for para in (body if isinstance(body, list) else [body]):
        out.append(parse_inline(para) if isinstance(para, str) else para)
    return out

def _blocks(blocks):
    out = []
    for b in (blocks or []):
        if isinstance(b, (list, tuple)):
            out.append(tuple(b)); continue
        t = b["type"]
        if t == "p":
            out.append(("p", parse_inline(b.get("text", ""))))
        elif t == "step":
            out.append(("step", b.get("badge", ""), b.get("lead", ""), b.get("text", "")))
        else:  # h, bullet, intro, warn, hint
            out.append((t, b.get("text", "")))
    return out

def _steps(steps):
    out = []
    for st in (steps or []):
        if isinstance(st, (list, tuple)): out.append(tuple(st))
        else: out.append((st.get("badge", ""), st.get("lead", ""), st.get("text", "")))
    return out

def _cards(cards):
    return [(c["label"], c["desc"]) if isinstance(c, dict) else tuple(c) for c in (cards or [])]

def _columns(cols):
    out = []
    for c in (cols or []):
        if isinstance(c, dict):
            out.append((c["head"], c.get("bullets", []), c.get("img")))
        else:
            out.append(tuple(c))
    return out

def _contact(contact):
    if not contact: return None
    return [(c["lead"], c["desc"]) if isinstance(c, dict) else tuple(c) for c in contact]

def _clean_opt(v):
    return v if v else None

def render_slide(sl):
    c = sl["component"]; g = sl.get
    if c == "cover":
        comp_cover(g("eyebrow", ""), g("title", ""), g("subtitle", ""), _clean_opt(g("image")))
    elif c == "content_card":
        comp_content_card(g("n"), g("eyebrow", ""), g("title", ""), _body(g("body", [])),
                          footnote=_clean_opt(g("footnote")), stats=g("stats") or None,
                          callout=_callout(g("callout")))
    elif c == "step_card":
        comp_step_card(g("n"), g("eyebrow", ""), g("title", ""), _steps(g("steps", [])),
                       intro=_clean_opt(g("intro")), callout=_callout(g("callout")),
                       image=_clean_opt(g("image")), note=_clean_opt(g("note")), hint=_clean_opt(g("hint")))
    elif c == "split":
        comp_split(g("n"), g("eyebrow", ""), g("title", ""), _blocks(g("blocks", [])),
                   _clean_opt(g("image")), callout=_callout(g("callout")))
    elif c == "two_col":
        comp_two_col(g("n"), g("eyebrow", ""), g("title", ""), _blocks(g("left", [])),
                     _blocks(g("right", [])), callout=_callout(g("callout")))
    elif c == "grid":
        comp_grid(g("n"), g("eyebrow", ""), g("title", ""), _cards(g("cards", [])),
                  cols=int(g("cols", 2)), intro=_clean_opt(g("intro")), badges=g("badges") or None,
                  callout=_callout(g("callout")), quote=_clean_opt(g("quote")))
    elif c == "grid3":
        comp_grid3(g("n"), g("eyebrow", ""), g("title", ""), _columns(g("columns", [])))
    elif c == "do_dont":
        comp_do_dont(g("n"), g("eyebrow", ""), g("title", ""), g("do", []), g("dont", []),
                     callout=_callout(g("callout")))
    elif c == "flow":
        comp_flow(g("n"), g("eyebrow", ""), g("title", ""), g("stages", []), g("caption", ""))
    elif c == "hero":
        comp_hero(g("n"), g("title", ""), subtitle=_clean_opt(g("subtitle")),
                  eyebrow=_clean_opt(g("eyebrow")), contact=_contact(g("contact")),
                  tagline=_clean_opt(g("tagline")), big=int(g("big", 60)))
    else:
        raise ValueError(f"Componente desconhecido: {c!r}")

def render_deck(deck, out=DEFAULT_OUT):
    global DATA
    _import_pptx()
    meta = deck.get("meta", {}) if isinstance(deck, dict) else {}
    slides = deck["slides"] if isinstance(deck, dict) else deck
    if meta.get("date"): DATA = meta["date"]
    _new_presentation()
    auto = 0
    for sl in slides:
        auto += 1
        sl.setdefault("n", auto)
        render_slide(sl)
    out = meta.get("output_file", out)
    if not os.path.isabs(out): out = os.path.join(BASE, out)
    prs.save(out)
    mb = os.path.getsize(out) / (1024 * 1024)
    print(f"OK -> {out}\nSlides: {len(prs.slides._sldIdLst)} | Tamanho: {mb:.2f} MB")
    return out

# ==================== DECK EMBUTIDO (Lei Lucas — 28 slides) ===================
DECK = {
    "meta": {
        "project": "Lei Lucas — Primeiros Socorros nas Escolas",
        "brand": "MAGMA Cursos",
        "date": DATA,
        "output_file": "lei_lucas_powerpoint_novo.pptx",
    },
    "slides": [
        {"component": "cover",
         "eyebrow": "TREINAMENTO · LEI Nº 13.722/2018",
         "title": "Lei Lucas: Primeiros Socorros nas Escolas",
         "subtitle": "Capacitação em Emergências Pré-Hospitalares — Diretrizes AHA 2025 | COFEN 805/2026 | SAMU/MS 2026",
         "image": "slide_01_image_1.png"},

        {"component": "content_card", "eyebrow": "Instrutor", "title": "Conhecendo nosso Instrutor",
         "body": ["**Enfermeiro** — Esp. em Enfermagem Neonatal e Pediatria; "
                  "Esp. Urgência e Emergência; Emergencista e Intensivista."],
         "stats": ["+01 mil alunos formados!", "08 anos ensinando pessoas a salvar vidas!"]},

        {"component": "content_card", "eyebrow": "Legislação", "title": "O que é a Lei Lucas?",
         "body": ["A **Lei nº 13.722/2018** obriga escolas e locais de recreação infantil a "
                  "capacitarem seus funcionários em primeiros socorros. Recebeu esse nome em "
                  "homenagem a Lucas Begalli, um menino de 10 anos que faleceu por engasgo em um "
                  "passeio escolar, sem que ninguém soubesse como agir."],
         "footnote": "A lei garante que crianças estejam mais protegidas — e reforça que qualquer "
                     "pessoa pode salvar uma vida com o conhecimento certo."},

        {"component": "content_card", "eyebrow": "Atualização", "title": "Resolução COFEN 805/2026",
         "body": ["Com a atualização da Resolução COFEN 805/2026, enfermeiros e técnicos de "
                  "enfermagem passaram a ter papel central na formação dessas equipes escolares, "
                  "garantindo que alunos, professores e funcionários estejam preparados para agir "
                  "nos primeiros minutos críticos de uma emergência."],
         "callout": {"kind": "alert", "title": None,
                     "desc": "Os primeiros 4 minutos são decisivos para salvar uma vida em parada cardiorrespiratória."}},

        {"component": "grid", "eyebrow": "Objetivos", "title": "Objetivos desta Aula", "cols": 2,
         "badges": ["1", "2", "3", "4"],
         "cards": [
             {"label": "Reconhecer PCR", "desc": "Identificar sinais de parada cardiorrespiratória e agir com segurança e rapidez."},
             {"label": "Realizar RCP", "desc": "Executar compressões e ventilações de alta qualidade conforme AHA 2025."},
             {"label": "Manobras de Engasgo", "desc": "Aplicar a manobra de Heimlich em adultos, crianças e bebês corretamente."},
             {"label": "AVC, Fraturas e mais", "desc": "Reconhecer AVC, queimaduras, fraturas e convulsões e saber como agir."}]},

        {"component": "two_col", "eyebrow": "PCR", "title": "O que é Parada Cardiorrespiratória (PCR)?",
         "left": [
             {"type": "h", "text": "Definição"},
             {"type": "p", "text": "A PCR é a cessação súbita das funções cardíacas e respiratórias, "
              "levando à interrupção do fluxo de sangue oxigenado para o cérebro e órgãos vitais. "
              "Sem intervenção imediata, o dano cerebral irreversível começa em 4 a 6 minutos."},
             {"type": "p", "text": "A PCR pode ocorrer em qualquer pessoa, em qualquer lugar — "
              "inclusive em crianças e jovens durante atividades físicas escolares."}],
         "right": [
             {"type": "h", "text": "Causas mais comuns"},
             {"type": "bullet", "text": "Infarto agudo do miocárdio"},
             {"type": "bullet", "text": "Choque elétrico"},
             {"type": "bullet", "text": "Afogamento"},
             {"type": "bullet", "text": "Trauma grave"},
             {"type": "bullet", "text": "Arritmias cardíacas"},
             {"type": "bullet", "text": "Engasgo com obstrução total"},
             {"type": "bullet", "text": "Overdose de medicamentos"}],
         "callout": {"kind": "alert", "title": None,
                     "desc": "Sem RCP imediata, a taxa de sobrevivência cai 10% a cada minuto!"}},

        {"component": "step_card", "eyebrow": "Reconhecimento", "title": "Como Reconhecer uma PCR?",
         "intro": "Antes de iniciar a RCP, confirme os sinais de PCR. Verifique os três elementos "
                  "abaixo rapidamente — todo esse processo deve durar no máximo 10 segundos.",
         "note": "Na dúvida, INICIE A RCP! O risco de não agir é sempre maior do que o de agir.",
         "steps": [
             {"badge": "1", "lead": "Inconsciência", "text": "A vítima não responde ao chamado, ao "
              "toque no ombro nem a estímulos dolorosos. Chame em voz alta: \"Você está bem?\""},
             {"badge": "2", "lead": "Ausência de Respiração Normal", "text": "Não há movimentos "
              "torácicos, ou a respiração é agônica (gasping) — sons irregulares e superficiais que "
              "não constituem respiração eficaz."},
             {"badge": "3", "lead": "Sem Pulso", "text": "Profissionais treinados verificam o pulso "
              "carotídeo por até 10 segundos. Para leigos, a ausência de resposta + ausência de "
              "respiração já indica PCR."}]},

        {"component": "step_card", "eyebrow": "Suporte Básico de Vida",
         "title": "Passo 1: Acione o SAMU — Ligue 192", "image": "slide_08_image_1.png",
         "steps": [
             {"badge": "01", "lead": "Garanta a Segurança", "text": "Verifique se o local é seguro "
              "para você e para a vítima. Não se exponha a riscos (elétricos, tráfego, estruturas instáveis)."},
             {"badge": "02", "lead": "Peça Ajuda em Voz Alta", "text": "Grite: \"Alguém me ajuda! "
              "Chame o SAMU — 192!\" Delegue a uma pessoa específica para ligar."},
             {"badge": "03", "lead": "Ligue 192 (SAMU)", "text": "Informe: localização exata, número "
              "de vítimas, condição da vítima e o que está sendo feito. Não desligue enquanto o operador orientar."},
             {"badge": "04", "lead": "Solicite o DEA", "text": "Se disponível, peça que alguém busque "
              "o Desfibrilador Externo Automático (DEA) enquanto você inicia a RCP."}]},

        {"component": "split", "eyebrow": "Passo 2 · Compressões",
         "title": "Compressões Torácicas de Alta Qualidade", "image": "slide_09_image_1.png",
         "blocks": [
             {"type": "intro", "text": "Como realizar compressões — Adulto"},
             {"type": "step", "badge": "", "lead": "Posicionamento", "text": "Coloque o calcanhar de "
              "uma mão no centro do peito (metade inferior do esterno). A outra mão por cima, dedos entrelaçados."},
             {"type": "step", "badge": "", "lead": "Profundidade", "text": "5 a 6 cm de profundidade. "
              "Braços estendidos, peso do corpo sobre as mãos."},
             {"type": "step", "badge": "", "lead": "Ritmo AHA 2025", "text": "100 a 120 compressões "
              "por minuto. Use a música \"Stayin' Alive\" (BeeGees) como referência de ritmo."},
             {"type": "step", "badge": "", "lead": "Recolhimento", "text": "Permita o recolhimento "
              "total do tórax entre as compressões. Não apoie o peso entre elas."}]},

        {"component": "step_card", "eyebrow": "Passo 3 · Ventilações", "title": "Ventilações — Relação 30:2",
         "intro": "Após 30 compressões, realize 2 ventilações de resgate. Cada ventilação deve durar "
                  "1 segundo, com volume suficiente para elevar visivelmente o tórax.",
         "image": "slide_11_image_2.png",
         "hint": "Se não souber ou não puder ventilar: faça RCP apenas com compressões (Hands-Only CPR), "
                 "conforme AHA 2025 para leigos.",
         "steps": [
             {"badge": "a", "lead": "Incline a Cabeça", "text": "Hiperextensão cervical: uma mão na "
              "testa, dois dedos sob o queixo. Abre a via aérea."},
             {"badge": "b", "lead": "Tampe o Nariz", "text": "Feche as narinas com o polegar e o "
              "indicador da mão que está na testa."},
             {"badge": "c", "lead": "Sopre Suavemente", "text": "Faça uma vedação com a boca e sopre "
              "por 1 segundo. Observe o tórax elevar."}]},

        {"component": "grid3", "eyebrow": "RCP Pediátrica", "title": "RCP Pediátrica — Crianças e Bebês",
         "columns": [
             {"head": "Bebê (0–12 meses)", "img": "slide_13_image_1.png", "bullets": [
                 "2 dedos no centro do peito (abaixo dos mamilos)", "Profundidade: 4 cm",
                 "Ritmo: 100–120/min | Relação 30:2", "Ventilação boca-a-boca-e-nariz"]},
             {"head": "Criança (1–8 anos)", "img": "slide_13_image_2.png", "bullets": [
                 "1 ou 2 mãos no centro do peito", "Profundidade: 5 cm (1/3 do diâmetro torácico)",
                 "Ritmo: 100–120/min | Relação 30:2", "Ventilação boca-a-boca"]},
             {"head": "Adulto (+8 anos)", "img": "slide_13_image_3.png", "bullets": [
                 "2 mãos entrelaçadas no centro do peito", "Profundidade: 5–6 cm",
                 "Ritmo: 100–120/min | Relação 30:2", "Ventilação boca-a-boca"]}]},

        {"component": "grid", "eyebrow": "Protocolo", "title": "Quando Parar a RCP?", "cols": 2,
         "intro": "A RCP deve ser mantida continuamente até que uma das condições abaixo ocorra. A "
                  "minimização das interrupções é um dos pilares das Diretrizes AHA 2025 — cada pausa "
                  "reduz a chance de sobrevivência.",
         "cards": [
             {"label": "SAMU Assume", "desc": "A equipe do SAMU ou outro profissional treinado chega e assume o atendimento."},
             {"label": "Retorno da Circulação", "desc": "A vítima retoma respiração espontânea, movimentos ou tosse — sinais de circulação espontânea."},
             {"label": "DEA Disponível", "desc": "O DEA é conectado — siga as instruções do dispositivo. Retome a RCP imediatamente após o choque."},
             {"label": "Exaustão Física", "desc": "O socorrista está incapaz de continuar. Se houver mais de uma pessoa, reveze a cada 2 minutos."}]},

        {"component": "flow", "eyebrow": "Fluxo de Decisão · BLS", "title": "Fluxograma BLS — AHA 2025",
         "stages": ["Reconhecer\nPCR", "Acionar\nSAMU (192)", "Iniciar RCP\n30:2", "Usar DEA\nAnalisar"],
         "caption": "Este fluxograma resume o algoritmo do Suporte Básico de Vida (BLS) para adultos "
                    "conforme as Diretrizes da American Heart Association 2025. A cadeia de sobrevivência "
                    "depende de reconhecimento precoce, RCP de alta qualidade e desfibrilação rápida."},

        {"component": "split", "eyebrow": "Engasgo · OVACE",
         "title": "Obstrução de Via Aérea por Corpo Estranho", "image": "slide_18_image_1.png",
         "blocks": [
             {"type": "h", "text": "Como Reconhecer o Engasgo Grave"},
             {"type": "bullet", "text": "Incapaz de falar ou emitir sons"},
             {"type": "bullet", "text": "Incapaz de tossir com força"},
             {"type": "bullet", "text": "Incapaz de respirar ou respiração muito difícil"},
             {"type": "bullet", "text": "Lábios e face ficando roxos (cianose)"},
             {"type": "bullet", "text": "Sinal universal: mãos no pescoço"}],
         "callout": {"kind": "info", "title": "Engasgo Leve — NÃO interfira!",
                     "desc": "Se a vítima consegue tossir com força, falar ou respirar, o engasgo é "
                             "leve. Encoraje-a a continuar tossindo — não bata nas costas nem faça compressões."}},

        {"component": "hero", "eyebrow": "OVACE", "title": "Tossir com força", "big": 70,
         "subtitle": "Lembre-se sempre de incentivar o paciente a continuar a tossir com força."},

        {"component": "hero", "eyebrow": "Manobras de Desobstrução", "title": "Técnicas", "big": 76,
         "subtitle": "Como agir diante do engasgo em diferentes situações."},

        {"component": "grid", "eyebrow": "Heimlich", "title": "Heimlich em Situações Especiais", "cols": 2,
         "cards": [
             {"label": "Gestante", "desc": "Compressões no tórax (não no abdômen). Punhos na metade inferior do esterno. Mesma lógica: 5 tapas nas costas + 5 compressões torácicas."},
             {"label": "Vítima Sentada", "desc": "Posicione-se atrás da cadeira. A técnica é a mesma — tapas nas costas e compressões abdominais com o punho abaixo do esterno."},
             {"label": "Engasgo Sozinho", "desc": "Faça as próprias compressões abdominais com o punho fechado. Ou jogue o abdômen contra uma superfície firme (encosto de cadeira, quina de mesa)."},
             {"label": "Vítima Inconsciente", "desc": "Deite a vítima no chão. Acione o SAMU (192) e inicie RCP. A cada ventilação, verifique se o objeto está visível para removê-lo com o dedo."}]},

        {"component": "split", "eyebrow": "Parada Respiratória",
         "title": "Parada Respiratória — Reconhecimento e Ação", "image": "slide_23_image_1.png",
         "blocks": [
             {"type": "h", "text": "O que é Parada Respiratória?"},
             {"type": "p", "text": "A parada respiratória ocorre quando a vítima para de respirar mas "
              "ainda tem pulso e batimento cardíaco. É diferente da PCR — o coração ainda bate, mas os "
              "pulmões pararam. Sem intervenção em minutos, evolui para PCR."},
             {"type": "h", "text": "Causas comuns"},
             {"type": "bullet", "text": "Afogamento"},
             {"type": "bullet", "text": "Choque elétrico"},
             {"type": "bullet", "text": "Obstrução parcial de via aérea"},
             {"type": "bullet", "text": "Overdose (drogas, medicamentos)"},
             {"type": "bullet", "text": "Convulsão prolongada"}]},

        {"component": "step_card", "eyebrow": "Ventilação de Resgate",
         "title": "Ventilação de Resgate — 1 sopro a cada 5 segundos",
         "intro": "Se a vítima tem pulso mas não respira, realize ventilações de resgate sem "
                  "compressões torácicas. A frequência recomendada pelas Diretrizes AHA 2025 é:",
         "callout": {"kind": "alert", "title": None,
                     "desc": "Mantenha a via aérea aberta durante todo o procedimento (hiperextensão "
                             "cervical ou mandíbula projetada em traumas)."},
         "steps": [
             {"badge": "1", "lead": "Sopro a cada 5 a 6 segundos", "text": "Equivalente a 10–12 ventilações por minuto em adultos."},
             {"badge": "1", "lead": "Sopro a cada 3 a 5 segundos", "text": "Em crianças e bebês — 12–20 ventilações por minuto."},
             {"badge": "10s", "lead": "Verificar pulso", "text": "A cada 2 minutos, reavalie pulso e respiração. Se perder o pulso, inicie RCP 30:2 imediatamente."}]},

        {"component": "step_card", "eyebrow": "Classificação e Conduta", "title": "Como Agir nas Queimaduras",
         "steps": [
             {"badge": "1", "lead": "1º Grau — Superficial", "text": "Vermelhidão e dor. Resfrie com água corrente por 20 minutos. Não aplique pasta de dente ou manteiga."},
             {"badge": "2", "lead": "2º Grau — Bolhas", "text": "Bolhas e dor intensa. Não estoure as bolhas. Cubra com pano limpo e procure atendimento médico."},
             {"badge": "3", "lead": "3º Grau — Profunda", "text": "Tecido escurecido, pouca dor (nervos destruídos). Ligue 192 imediatamente. Não toque na área queimada."}]},

        {"component": "do_dont", "eyebrow": "Conduta", "title": "Queimaduras — O que Fazer e o que NÃO Fazer",
         "do": ["Resfrie imediatamente com água fria corrente por 10 a 20 minutos",
                "Retire roupas e acessórios ao redor da queimadura (exceto se grudados)",
                "Cubra com pano limpo ou curativo não aderente",
                "Ligue 192 (SAMU) em queimaduras extensas, profundas ou em face/mãos/genitais"],
         "dont": ["Não use gelo — causa vasoconstrição e aprofunda a lesão",
                  "Não fure bolhas — proteção natural contra infecção",
                  "Não aplique pasta de dente, manteiga, óleo ou qualquer produto caseiro",
                  "Não use água gelada ou muito quente",
                  "Não remova roupa grudada na pele"]},

        {"component": "do_dont", "eyebrow": "Convulsões", "title": "Convulsões — Proteger, Não Conter",
         "do": ["Proteja a cabeça com algo macio (roupa dobrada)",
                "Afaste objetos que possam machucar",
                "Após a crise: posição lateral de segurança (PLS)",
                "Cronometre a duração da crise",
                "Ligue 192 se durar +5 min ou se for a primeira vez"],
         "dont": ["Não segure ou imobilize a vítima à força",
                  "Não coloque nada na boca — engolir a língua é um mito",
                  "Não ofereça água ou remédios durante a crise",
                  "Não deixe a vítima sozinha"],
         "callout": {"kind": "info", "title": None,
                     "desc": "A Posição Lateral de Segurança (PLS) evita que a vítima se afogue com secreções após a crise."}},

        {"component": "grid", "eyebrow": "AVC", "title": "AVC — Acidente Vascular Cerebral", "cols": 2,
         "badges": ["F", "A", "S", "T"],
         "intro": "O AVC ocorre quando o suprimento de sangue ao cérebro é interrompido ou reduzido. "
                  "Cada minuto sem tratamento, 1,9 milhão de neurônios morrem. Reconhecer precocemente "
                  "salva vida e qualidade de vida.",
         "cards": [
             {"label": "F — Face", "desc": "Peça para sorrir. Um lado da boca caiu? Assimetria facial é sinal claro de AVC."},
             {"label": "A — Braço", "desc": "Levante os dois braços. Um cai involuntariamente? Fraqueza unilateral é sinal de alarme."},
             {"label": "S — Fala", "desc": "Peça para repetir uma frase simples. Fala arrastada, confusa ou incapacidade de falar = emergência."},
             {"label": "T — Tempo", "desc": "Anote o horário dos primeiros sintomas e ligue 192 imediatamente. Tempo é cérebro!"}]},

        {"component": "step_card", "eyebrow": "Como Identificar", "title": "Método SAMU — Sinais de AVC",
         "steps": [
             {"badge": "1", "lead": "Sorriso", "text": "Peça para sorrir. Um lado do rosto cai ou fica paralisado?"},
             {"badge": "2", "lead": "Abraço", "text": "Peça para levantar os dois braços. Um braço cai ou não se move?"},
             {"badge": "3", "lead": "Música", "text": "Peça para repetir uma frase simples. A fala está enrolada ou confusa?"},
             {"badge": "4", "lead": "Urgente", "text": "Se qualquer sinal estiver presente, ligue 192 (SAMU) imediatamente!"}]},

        {"component": "grid", "eyebrow": "Conclusão", "title": "Cada Segundo Conta", "cols": 3,
         "quote": "A diferença entre uma vida salva e uma vida perdida pode ser apenas o conhecimento de quem estava presente.",
         "callout": {"kind": "success", "title": None,
                     "desc": "Você agora sabe como agir em PCR, engasgo, queimadura, fratura, convulsão, AVC e parada respiratória. Parabéns!"},
         "cards": [
             {"label": "RCP Salva Vidas", "desc": "Compressões de alta qualidade (100–120/min, 5–6 cm, relação 30:2) são a principal arma contra a PCR antes do SAMU chegar."},
             {"label": "A Escola é um Ponto Chave", "desc": "A Lei Lucas garante que profissionais de educação estejam preparados. Cada escola treinada é uma comunidade mais segura."},
             {"label": "192 — Sempre", "desc": "Em qualquer emergência, ligue 192 (SAMU). Não tente transportar sozinho. Aguarde orientação profissional."}]},

        {"component": "grid", "eyebrow": "Base Legal", "title": "Referências e Base Legal", "cols": 2,
         "callout": {"kind": "tip", "title": None,
                     "desc": "Coordenação de Enfermagem — CISBAF/SAMU  |  Material atualizado 2026"},
         "cards": [
             {"label": "American Heart Association (AHA)", "desc": "Guidelines for CPR and Emergency Cardiovascular Care — 2025 Update: RCP de alta qualidade, algoritmos BLS adulto e pediátrico, minimização de interrupções, Hands-Only CPR."},
             {"label": "Resolução COFEN 805/2026", "desc": "Dispõe sobre a atuação do enfermeiro e técnico de enfermagem na capacitação e supervisão de primeiros socorros em instituições de ensino, no âmbito da Lei Lucas."},
             {"label": "SAMU / Ministério da Saúde — 2026", "desc": "Protocolos de Atendimento Pré-Hospitalar do SAMU 192, versão 2026. Ministério da Saúde — Secretaria de Atenção Especializada à Saúde."},
             {"label": "Lei nº 13.722/2018 — Lei Lucas", "desc": "Torna obrigatória a capacitação em primeiros socorros de professores e funcionários de estabelecimentos de educação básica. Publicada no DOU em 04/10/2018."}]},

        {"component": "hero", "eyebrow": "Encerramento", "title": "Salvar vidas começa com conhecimento", "big": 50,
         "subtitle": "Agora você conhece os principais protocolos de primeiros socorros baseados nas "
                     "diretrizes mais atuais. Compartilhe esse conhecimento — ele pode fazer a diferença "
                     "quando cada segundo conta.",
         "contact": [
             {"lead": "Magma Cursos", "desc": "Capacitação prática em primeiros socorros"},
             {"lead": "Enf. João Paulo Bello", "desc": "Dúvidas? Entre em contato!"}]},

        {"component": "hero", "title": "Obrigado!", "tagline": "Educar também é cuidar.", "big": 96},
    ],
}

# ==================== CLI ====================================================
def _load(path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def main(argv):
    args = [a for a in argv if not a.startswith("--")]
    if "--dump" in argv:
        out = args[0] if args else os.path.join(BASE, "deck.json")
        with open(out, "w", encoding="utf-8") as f:
            json.dump(DECK, f, ensure_ascii=False, indent=2)
        print(f"deck escrito em {out}  ({len(DECK['slides'])} slides)")
        return
    deck = _load(args[0]) if args else DECK
    out = args[1] if len(args) > 1 else DEFAULT_OUT
    render_deck(deck, out)

if __name__ == "__main__":
    main(sys.argv[1:])
